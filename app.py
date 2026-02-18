import streamlit as st
import pandas as pd
import os
import requests
import re
import random
import io
from groq import Groq
from pptx import Presentation # Slayd uchun

# --- 1. SOZLAMALAR ---
try:
    MAKTAB_NOMI = "1-sonli umumta'lim maktabi"
    DIREKTOR_FIO = "Mahmudov Matyoqub Narzulloyevich"
    TO_GRI_PAROL = "informatika2024"
    MONITORING_KODI = "admin777"
    
    BOT_TOKEN = st.secrets["TELEGRAM_BOT_TOKEN"]
    GURUH_ID = st.secrets["TELEGRAM_GURUH_ID"]
    GROQ_API_KEY = st.secrets["GROQ_API_KEY"]
    
    client = Groq(api_key=GROQ_API_KEY)
except Exception:
    st.error("⚠️ Secrets ma'lumotlarida xatolik bor!")
    st.stop()

# --- 2. BAZANI YUKLASH ---
@st.cache_data
def yuklash():
    files = [f for f in os.listdir('.') if f.lower().endswith(('.xlsx', '.xls'))]
    all_sheets = {}
    for f in files:
        try:
            sheets = pd.read_excel(f, sheet_name=None, dtype=str)
            for name, df in sheets.items():
                if not df.empty:
                    df.columns = [str(c).strip().lower() for c in df.columns]
                    all_sheets[f"{f} | {name}"] = df
        except: continue
    return all_sheets

sheets_baza = yuklash()

# --- 3. SLAYD (PPTX) YARATISH ---
def pptx_yarat(df, sarlavha):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = f"Hisobot: {sarlavha.upper()}"
    slide.placeholders[1].text = df.to_string(index=False)[:1000]
    ppt_io = io.BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)
    return ppt_io

# --- 4. DIZAYN VA KIRISH ---
st.set_page_config(page_title=MAKTAB_NOMI, layout="wide")

if "authenticated" not in st.session_state:
    st.title(f"🏫 {MAKTAB_NOMI}")
    p_in = st.text_input("Kirish paroli:", type="password", key="main_auth_key")
    if st.button("Kirish", key="main_auth_btn"):
        if p_in == TO_GRI_PAROL:
            st.session_state.authenticated = True
            st.rerun()
        else: st.error("Xato!")
    st.stop()

# --- 5. SIDEBAR ---
with st.sidebar:
    st.title(f"🏛 {MAKTAB_NOMI}")
    st.write(f"👤 **Direktor:** {DIREKTOR_FIO}")
    menu = st.radio("Bo'limni tanlang:", ["🤖 AI Muloqot", "📊 Jurnal Monitoringi"], key="nav_menu")
    st.divider()
    st.info("💡 Bilim - najotdir.")

# --- 6. AI MULOQOT (FAROSATLI VA SAMIMIY) ---
if menu == "🤖 AI Muloqot":
    st.title("🤖 Maktab AI yordamchisi")
    if "messages" not in st.session_state: st.session_state.messages = []
    
    for msg in st.session_state.messages:
        with st.chat_message(msg["role"]): st.markdown(msg["content"])

    if savol := st.chat_input("Ism yozing yoki savol so'rang...", key="chat_input_unique"):
        st.session_state.messages.append({"role": "user", "content": savol})
        with st.chat_message("user"): st.markdown(savol)
        
        with st.chat_message("assistant"):
            q_low = savol.lower().strip()
            
            # A) ODAM OMILI: SAMIMIY JAVOBLAR
            samimiy_javob = None
            if any(x in q_low for x in ["rahmat", "tashakkur", "bor bo'ling", "baraka toping"]):
                samimiy_javob = "Arzimaydi! Sizdan ham Alloh rozi bo'lsin. Maktabimiz ravoji uchun xizmat qilishdan xursandman. 😊"
            elif any(x in q_low for x in ["salom", "assalom", "qalaysiz", "yaxshimisiz"]):
                samimiy_javob = "Vaalaykum assalom! Men yaxshiman, rahmat. O'zingiz yaxshi yuribsizmi? Sizga qanday yordam bera olaman? 🏛"
            elif any(x in q_low for x in ["xayr", "xush qoling", "ertagacha", "ko'rishguncha"]):
                samimiy_javob = "Xayr, sog'-salomat bo'ling! Ishlaringizda omad tilayman. 👋"
            elif any(x in q_low for x in ["zo'r", "ajoyib", "gap yo'q", "yaxshi ishlayapsan"]):
                samimiy_javob = "Katta rahmat! Maqtovlaringizdan ruhlandim. Sizga yordam berish men uchun sharaf! 🌟"

            if samimiy_javob:
                st.markdown(samimiy_javob)
                st.session_state.messages.append({"role": "assistant", "content": samimiy_javob})
            
            # B) BAZADAN QIDIRUV (Agar samimiy suhbat bo'lmasa)
            else:
                topildi = False
                clean_word = q_low.replace("ro'yxati", "").replace("top", "").replace("o'qituvchilar", "").strip()
                
                if len(clean_word) >= 3:
                    for key, df in sheets_baza.items():
                        mask = df.apply(lambda r: r.astype(str).str.contains(clean_word, case=False, na=False).any(), axis=1)
                        res_df = df[mask]
                        if not res_df.empty:
                            st.success(f"🔍 '{clean_word}' bo'yicha ma'lumot:")
                            st.dataframe(res_df, use_container_width=True)
                            
                            # PPTX Slayd yuklash tugmasi
                            ppt_data = pptx_yarat(res_df, savol)
                            st.download_button(
                                label="📂 Slaydni PPTX shaklida yuklash",
                                data=ppt_data,
                                file_name=f"{clean_word}_hisobot.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                key=f"dl_{random.randint(0,999)}"
                            )
                            topildi = True
                            break
                
                # C) AGAR BAZADA BO'LMASA -> GROQ AI
                if not topildi:
                    try:
                        res = client.chat.completions.create(
                            messages=[{"role":"system","content":f"Sen {MAKTAB_NOMI} AI yordamchisisan. O'zbekona lutf bilan, samimiy va madaniyatli gaplash."},
                                     {"role":"user","content":savol}],
                            model="llama-3.3-70b-versatile"
                        )
                        msg_text = res.choices[0].message.content
                        st.markdown(msg_text)
                        st.session_state.messages.append({"role": "assistant", "content": msg_text})
                    except:
                        st.error("AI hozirda band.")

# --- 7. MONITORING (SENING MANTIQING - ASL HOLIDA) ---
elif menu == "📊 Jurnal Monitoringi":
    st.title("📊 Jurnal Monitoringi")
    if "m_auth" not in st.session_state: st.session_state.m_auth = False
    if not st.session_state.m_auth:
        m_input = st.text_input("Monitoring kodi:", type="password", key="mon_input")
        if st.button("Kirish", key="mon_btn"):
            if m_input == MONITORING_KODI: 
                st.session_state.m_auth = True
                st.rerun()
            else: st.error("Kod xato!")
        st.stop()
    
    j_fayl = st.file_uploader("Excel faylni yuklang", type=['xlsx', 'xls', 'html'], key="uploader")
    if j_fayl:
        try:
            # Faylni o'qish (engine'lar bilan)
            try: df_j = pd.read_excel(j_fayl, engine='openpyxl')
            except:
                try: 
                    j_fayl.seek(0)
                    df_j = pd.read_excel(j_fayl, engine='xlrd')
                except:
                    j_fayl.seek(0)
                    df_j = pd.read_html(j_fayl, header=0)[0]

            df_j.columns = [str(c).replace('\n', ' ').strip() for c in df_j.columns]
            kamchiliklar = []
            if len(df_j.columns) >= 6:
                for _, row in df_j.iterrows():
                    name, val = str(row.iloc[0]), str(row.iloc[5])
                    if any(x in name.lower() for x in ["tuman", "muassasa", "f.i.sh"]): continue
                    nums = re.findall(r'(\d+)', val)
                    if len(nums) >= 2 and int(nums[0]) < int(nums[1]):
                        kamchiliklar.append(f"❌ **{name}**: {int(nums[1])-int(nums[0])} ta chala ({val})")
                
                st.dataframe(df_j, use_container_width=True)
                xabar_text = "✅ Hammasi to'liq!" if not kamchiliklar else "⚠️ **Kamchiliklar:**\n\n" + "\n".join(kamchiliklar)
                st.warning(xabar_text) if kamchiliklar else st.success(xabar_text)
                
                if st.button("📢 Telegramga yuborish", key="tg_btn"):
                    requests.post(f"https://api.telegram.org/bot{BOT_TOKEN}/sendMessage", 
                                 json={"chat_id": GURUH_ID, "text": f"📊 <b>Monitoring</b>\n\n{xabar_text}", "parse_mode": "HTML"})
                    st.success("✅ Telegramga yuborildi!")
        except Exception as e: st.error(f"Xato: {e}")
